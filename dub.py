
# coding: utf-8

# 还是不甘心, 想自动加入音频试试

# In[1]:


from pptx import Presentation
from generate_voice import tts
import os,sys
import re
import subprocess
from pptx.util import Inches

from add_voice import get_notes_text, get_notes_voice, save_notes_voice


# In[2]:


ppt_filename="data/test.pptx"


# In[3]:


ppt_path=os.path.dirname(ppt_filename)
temp_path=os.path.join(ppt_path, "temp")
try:
    os.mkdir(temp_path)
except:
    pass

notes_text=get_notes_text(ppt_filename)
notes_voice=get_notes_voice(notes_text)
movie_list=save_notes_voice(notes_voice, temp_path)


# In[4]:


prs = Presentation(ppt_filename)
left = top = Inches(0.0)
width = height = Inches(1.0)

for index, slide in enumerate(prs.slides):
    notes_slide = slide.notes_slide
    shapes = slide.shapes
    movie = shapes.add_movie(movie_list[index], 
                             left , top , width , height, 
                             poster_frame_image=None, 
                             mime_type='video/unknown')

prs.save(os.path.join(ppt_path, "output.pptx"))


# In[5]:


# for file in movie_list.values():
#     os.remove(file)
# os.rmdir(temp_path)

