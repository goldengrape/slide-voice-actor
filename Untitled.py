
# coding: utf-8

# 还是不甘心, 想自动加入音频试试

# In[65]:


from pptx import Presentation
from generate_voice import tts
import os
import re
import subprocess
from pptx.util import Inches

from add_voice import get_notes_text, get_notes_voice, save_notes_voice


# In[66]:


ppt_filename="data/test.pptx"


movie_file=["data/test/temp_000.mp3","data/test/temp_001.mp3"]
prs = Presentation(ppt_filename)
left = top = Inches(0.0)
width = height = Inches(1.0)

for index, slide in enumerate(prs.slides):
    notes_slide = slide.notes_slide
    shapes = slide.shapes
    movie = shapes.add_movie(movie_file[index], 
                             left , top , width , height, 
                             poster_frame_image=None, 
                             mime_type='video/unknown')

prs.save('data/test2.pptx')


# In[44]:


prs2=Presentation('data/test2.pptx')
shapes=prs2.slides[0].shapes


# In[58]:


shapes[-1].shape_type


# In[35]:


dir(shapes[-1])

