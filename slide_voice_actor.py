
# coding: utf-8

# # 使用mac自带的TTS引擎为PPT配音
# 

# ## 初始化

# In[1]:


import sys, os , subprocess, time
from pptx import Presentation
from  AppKit import NSSpeechSynthesizer
import Foundation
from pptx.util import Inches


# In[2]:


def read_pptx(ppt_filename):
    prs = Presentation(ppt_filename)
    return prs

def init_tts(rate=200):
    nssp = NSSpeechSynthesizer
    ve = nssp.alloc().init()
    ve.setRate_(rate)
    return ve


# ## 读取ppt中每页的注释

# In[3]:


def get_notes_text(slide):
    notes_slide = slide.notes_slide
    note = notes_slide.notes_text_frame.text
    return note


# ## 将每页注释用tts转换成音频文件

# In[4]:


def save_tts(ve, TEXT, filename):
    url = Foundation.NSURL.fileURLWithPath_(filename)
    s=ve.startSpeakingString_toURL_(TEXT,url)    
    if not(s):
        print("TTS failed") 
    return filename


# In[5]:


def save_notes_voice(ve, text, page_number):
    voice_filename= "temp_tts_{:3d}.aiff".format(page_number)
    voice_filename= save_tts(ve, text,voice_filename )
    return voice_filename


# In[6]:


def convert_aiff_to_mp3(filename):
    # ffmpeg -i myinput.aif -f mp3 -acodec libmp3lame -ab 320000 -ar 44100 myoutput.mp3
    mp3_filename=filename+".mp3"
    subprocess.call(["ffmpeg",  
                     "-i", filename,
                     "-f", "mp3",
                     "-acodec", "libmp3lame",
                     "-ab", "128000",
                     "-ar", "44100",
                     mp3_filename])
    print("made", mp3_filename)
    return mp3_filename


# ## 将每个音频文件插入到ppt页面中

# In[7]:


def insert_voice(voice_filename, slide):
    left = top = Inches(0.0)
    width = height = Inches(1.0)

    shapes = slide.shapes
    movie = shapes.add_movie(voice_filename, 
                                 left , top , width , height, 
                                 poster_frame_image=None, 
                                 mime_type='video/unknown')


# ## 清理掉临时文件

# In[8]:


def clean_temp(voice_filename):
    os.remove(voice_filename)
    


# # 包装

# In[9]:


def main(ppt_filename, output_filename):
    ve=init_tts(rate=200)
    prs=read_pptx(ppt_filename)
    for index, slide in enumerate(prs.slides): 
        note=get_notes_text(slide)
        voice_filename=save_notes_voice(ve, note, index)
#         voice_filename=convert_aiff_to_mp3(voice_filename)
        time.sleep(3)
        insert_voice(voice_filename, slide)
        print("Slide No. {}".format(index))
#         clean_temp(voice_filename)
    prs.save(output_filename)
    print("save to ",output_filename)


# In[10]:


if __name__=="__main__":
    if len(sys.argv)>=3 :
        ppt_filename=sys.argv[1]
        output_filename=sys.argv[2]
    elif len(sys.argv)==2 :
        ppt_filename=sys.argv[1]
        ppt_path=os.path.dirname(ppt_filename)
        output_filename=os.path.join(ppt_path, "output.pptx")
    else:
        print("Error, I need input filename")
        
    ppt_filename="data/test.pptx"
    ppt_path=os.path.dirname(ppt_filename)
    output_filename=os.path.join(ppt_path, "output.pptx")
    
    main(ppt_filename, output_filename)

