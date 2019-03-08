
# coding: utf-8

# 读取ppt, 将ppt中的注释文字提取出来

# In[1]:


from pptx import Presentation
from generate_voice import tts
import os
import re
import subprocess


# In[2]:


if __name__=="__main__":
    ppt_filename="data/test.pptx"


# In[3]:


def get_notes_text(ppt_filename):
    prs = Presentation(ppt_filename)
    dict={}
    for index, slide in enumerate(prs.slides):
        notes_slide = slide.notes_slide
        text = notes_slide.notes_text_frame.text
        dict[index]=text
    return dict


# 将每个注释转换成语音

# In[4]:


def get_notes_voice(notes_dict):
    voice_dict={}
    for index, text in notes_dict.items():
        voice_dict[index]=tts(text)
    return voice_dict


# 将语音分别保存为音频mp3

# In[5]:


if __name__=="__main__":

    notes_dict=get_notes_text(ppt_filename)
    voice_dict=get_notes_voice(notes_dict)


# In[6]:


def save_notes_voice(voice_dict,images_path):
    temp_file_dict={}
    for index, voice_data in voice_dict.items():
        temp_filename=os.path.join(images_path,"temp_{:03d}.mp3".format(index))
        with open(temp_filename, 'wb') as outfile:
            outfile.write(voice_data)
#             print("save ", temp_filename)
        temp_file_dict[index]=temp_filename
    return temp_file_dict


# 将每页导出成图像, 然后把每张图像与每个音频做成一个视频

# In[7]:


def build_slide_video(images_path):
    image_filename_list = [f for f in os.listdir(images_path) if f.endswith('.png')]
    audio_filename_list = save_notes_voice(voice_dict,images_path)
    video_filename_list = ["slide_{:03d}.mp4".format(index)
                           for index in range(len(image_filename_list))]
    for filename in image_filename_list:
        (prefix, index_str,ends)=re.findall("([\u4e00-\u9fa5A-Z]+)(\d+)(.png)",filename)[0]
        index=int(index_str)-1 # 注意python从0开始, ppt导出图片从1开始

        image_filename=os.path.join(images_path,filename)
        audio_filename=audio_filename_list[index]
        video_filename=os.path.join(images_path,video_filename_list[index])

        ffmpeg_str="ffmpeg -loop 1 -i {} -i {} -c:v libx264 -c:a copy -shortest {}".format(
            image_filename,audio_filename,video_filename
        )
        subprocess.call(["ffmpeg", "-loop", "1", 
                         "-i", image_filename,
                         "-i", audio_filename,
                         "-c:v", "libx264", "-c:a", "copy", "-shortest",
                         video_filename])
#         print("save ", video_filename )
    # 记录产生的文件名
    video_list=os.path.join(images_path,"video_file_list.txt")
    with open(video_list, "w") as f:
        for v in video_filename_list:
            f.write("file "+"'{}'\n".format(v)) # 注意ffmpeg在从txt里调用文件名时, 似乎是用了相对路径
    
    return video_list


# 把所有视频连起来

# In[8]:


def generate_slides_with_voice(ppt_filename):
    images_path=ppt_filename.split(".")[0]
    video_list=build_slide_video(images_path)
    # ffmpeg -f concat -i mylist.txt -c copy output
    subprocess.call(["ffmpeg", "-f", "concat",
                     "-i", video_list,
                     "-c", "copy",
                     ppt_filename+".mp4"])


# In[9]:


if __name__=="__main__":
    generate_slides_with_voice(ppt_filename)

